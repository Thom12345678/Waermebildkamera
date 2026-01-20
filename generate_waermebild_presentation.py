import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

def create_presentation(output_path, author, date):
    prs = Presentation()

    # Folien-Inhalte für eine allgemeine Präsentation über Wärmebildkameras
    slides_data = [
        ["Grundlagen der Wärmebildtechnik", 
         f"Alles über Infrarot und Thermografie\n\nErstellt von: {author}\nDatum: {date}"],
        
        ["Was ist eine Wärmebildkamera?", 
         "• Optisches Gerät zur Messung von Infrarotstrahlung\n• Wandelt Wärmesignale in ein sichtbares Bild um (Falschfarbenbild)\n• Erfasst Oberflächen-Temperaturen berührungslos"],
        
        ["Das Prinzip der Infrarotstrahlung", 
         "• Jeder Körper über -273,15 °C sendet Strahlung aus\n• Je wärmer, desto intensiver die Strahlung\n• Wellenlängenbereich: ca. 8 bis 14 µm"],
        
        ["Wichtige Begriffe der Thermografie", 
         "• Emissionsgrad: Wie gut strahlt ein Material Wärme ab?\n• Reflexion: Vorsicht bei glänzenden Oberflächen!\n• Transmission: Durch welche Stoffe geht die Strahlung (z.B. Germanium-Linsen)?"],
        
        ["Einsatz in der Elektrotechnik", 
         "• Prüfung von Schaltschränken unter Last\n• Erkennen von lockeren Verbindungen und Überlastungen\n• Brandschutz durch Hotspot-Erkennung"],
        
        ["Einsatz im Bauwesen & Industrie", 
         "• Gebäude-Thermografie (Wärmebrücken, Isolation)\n• Lokalisierung von Leckagen in Rohrleitungen\n• Überwachung von Maschinenlagern und Motoren"],
        
        ["Einsatz bei Feuerwehr & Rettung", 
         "• Personensuche in verrauchten Räumen\n• Glutnestsuche nach Bränden\n• Kitzrettung in der Landwirtschaft"],
        
        ["Vorteile der Technik", 
         "• Messung im laufenden Betrieb möglich\n• Großflächige Überwachung in Echtzeit\n• Präventive Wartung spart hohe Reparaturkosten"],
        
        ["Grenzen und Fehlerquellen", 
         "• Glas ist für IR-Kameras undurchsichtig\n• Starke Reflexionen auf polierten Metallen\n• Nur Oberflächentemperaturen messbar"],
        
        ["Fazit & Ausblick", 
         "• Unverzichtbares Werkzeug für Technik und Sicherheit\n• Kameras werden immer kompakter und günstiger\n• Integration in Smartphones und Drohnen"]
    ]

    for title_text, body_text in slides_data:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        
        content = slide.placeholders[1]
        content.text = body_text

    prs.save(output_path)
    print(f"Präsentation erfolgreich erstellt: {output_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", default="presentations/Waermebildkamera_6min_Praesentation.pptx")
    parser.add_argument("--author", default="Thom12345678")
    parser.add_argument("--date", default=datetime.now().strftime("%d.%m.%Y"))
    args = parser.parse_args()

    create_presentation(args.output, args.author, args.date)